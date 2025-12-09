#!/usr/bin/env python3
"""
Create professional FinTech-style pitch deck for FE571 Final Presentation
Uses Stevens Institute branding and Helvetica Neue font styling
"""

import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Stevens Institute Colors
STEVENS_RED = RGBColor(0x9D, 0x15, 0x35)
STEVENS_GRAY = RGBColor(0x94, 0x95, 0x94)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# Font settings (Helvetica Neue - will fallback to Helvetica/Arial if not available)
FONT_PRIMARY = 'Helvetica Neue'
FONT_FALLBACK = 'Arial'

# Charts directory
CHARTS_DIR = '/Users/akbarpathan/Desktop/Dev/QuantHFStrat/backtesting/charts/'

def set_slide_background(slide, color=WHITE):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_speaker_notes(slide, notes):
    """Add speaker notes to a slide (list of bullet points)"""
    if notes:
        notes_text = "\n".join(f"• {note}" for note in notes)
        slide.notes_slide.notes_text_frame.text = notes_text

def add_title_slide(prs, title, subtitle, authors, notes=None):
    """Add title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(24)
    p.font.color.rgb = STEVENS_RED
    p.alignment = PP_ALIGN.CENTER

    # Authors
    auth_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    tf = auth_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FE571 | Professor Anshul Sharma | Group 7"
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(16)
    p.font.color.rgb = STEVENS_GRAY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = authors
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(14)
    p.font.color.rgb = STEVENS_GRAY
    p.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide, notes)
    return slide

def add_content_slide(prs, title, bullets, one_liner=None, notes=None):
    """Add a content slide with title, bullets, and optional one-liner"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Red accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STEVENS_RED
    line.line.fill.background()

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    # One-liner at bottom
    if one_liner:
        liner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
        tf = liner_box.text_frame
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = STEVENS_GRAY

    add_speaker_notes(slide, notes)
    return slide

def add_table_slide(prs, title, headers, rows, one_liner=None, notes=None):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Red accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STEVENS_RED
    line.line.fill.background()

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(9)
    table_height = Inches(0.4 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.4), table_width, table_height).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = STEVENS_RED
        para = cell.text_frame.paragraphs[0]
        para.font.name = FONT_PRIMARY
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            para = cell.text_frame.paragraphs[0]
            para.font.name = FONT_PRIMARY
            para.font.size = Pt(12)
            para.font.color.rgb = DARK_GRAY
            para.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # One-liner
    if one_liner:
        liner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
        tf = liner_box.text_frame
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = STEVENS_GRAY

    add_speaker_notes(slide, notes)
    return slide

def add_chart_placeholder_slide(prs, title, chart_name, one_liner=None):
    """Add a slide with placeholder for chart"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Red accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STEVENS_RED
    line.line.fill.background()

    # Chart placeholder box
    chart_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.4), Inches(8.5), Inches(3.5))
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    chart_box.line.color.rgb = STEVENS_GRAY

    # Placeholder text
    text_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(6), Inches(1))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"[INSERT: {chart_name}]"
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(18)
    p.font.color.rgb = STEVENS_GRAY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Export from v5 notebook: backtesting/charts/"
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(12)
    p.font.color.rgb = STEVENS_GRAY
    p.alignment = PP_ALIGN.CENTER

    # One-liner
    if one_liner:
        liner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
        tf = liner_box.text_frame
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = STEVENS_GRAY

    return slide

def add_chart_image_slide(prs, title, image_filename, one_liner=None, notes=None):
    """Add a slide with an actual chart image embedded"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Red accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STEVENS_RED
    line.line.fill.background()

    # Image
    image_path = os.path.join(CHARTS_DIR, image_filename)

    if os.path.exists(image_path):
        # Get image dimensions
        with Image.open(image_path) as img:
            img_width, img_height = img.size

        aspect_ratio = img_width / img_height

        # Available space (slide is 10" x 5.625")
        max_width = 8.5  # inches (leave 0.75" margins)
        max_height = 3.5  # inches (leave room for title + one-liner)

        # Calculate size to fit while preserving aspect ratio
        if aspect_ratio > (max_width / max_height):
            # Width-constrained (wide images)
            width = max_width
            height = width / aspect_ratio
        else:
            # Height-constrained (square-ish images)
            height = max_height
            width = height * aspect_ratio

        # Center horizontally
        left = (10 - width) / 2  # 10" slide width
        top = 1.3  # Below title and accent line

        slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                  width=Inches(width), height=Inches(height))
    else:
        # Fallback placeholder if image not found
        chart_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.4), Inches(8.5), Inches(3.5))
        chart_box.fill.solid()
        chart_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        chart_box.line.color.rgb = STEVENS_GRAY

        text_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(6), Inches(1))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[IMAGE NOT FOUND: {image_filename}]"
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(18)
        p.font.color.rgb = STEVENS_GRAY
        p.alignment = PP_ALIGN.CENTER

    # One-liner at bottom
    if one_liner:
        liner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
        tf = liner_box.text_frame
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = STEVENS_GRAY

    add_speaker_notes(slide, notes)
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, one_liner=None, notes=None):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Red accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STEVENS_RED
    line.line.fill.background()

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = STEVENS_RED

    # Left column items
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(3))
    tf = left_box.text_frame
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = STEVENS_RED

    # Right column items
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4.2), Inches(3))
    tf = right_box.text_frame
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # One-liner
    if one_liner:
        liner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
        tf = liner_box.text_frame
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.name = FONT_PRIMARY
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = STEVENS_GRAY

    add_speaker_notes(slide, notes)
    return slide

def add_questions_slide(prs, notes=None):
    """Add questions slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Questions text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # GitHub link
    link_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    tf = link_box.text_frame
    p = tf.paragraphs[0]
    p.text = "github.com/Ayan-Mahmood/QuantHFStrat"
    p.font.name = FONT_PRIMARY
    p.font.size = Pt(16)
    p.font.color.rgb = STEVENS_RED
    p.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide, notes)
    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # SLIDE 1: Title
    add_title_slide(
        prs,
        "Statistical Basket Pairs Trading Strategy",
        "Volatility Dispersion Mean-Reversion",
        "Scott Henriquez, Nakul Jadeja, Ayan Mahmood, Akbar Pathan",
        notes=[
            "This project evolved from a simple pairs trading idea into a full quantitative research exercise",
            "We'll walk through our methodology, what worked, what didn't, and the lessons learned along the way",
            "The code is fully reproducible on GitHub — we welcome questions at the end"
        ]
    )

    # SLIDE 2: Executive Summary
    add_content_slide(
        prs,
        "Executive Summary",
        [
            "Long/short volatility spread strategy across 4 sector pairs",
            "Market-neutral design with low SPY correlation (-0.07)",
            "ML filter tested to improve signal quality",
            "Result: Solid framework, inconsistent alpha — best as portfolio hedge"
        ],
        notes=[
            "We're being upfront: this strategy doesn't generate consistent alpha, but that's actually an important finding",
            "The real value is in the process — building a rigorous backtesting framework with proper bias controls",
            "Market neutrality is valuable even without alpha — it provides diversification benefits in a portfolio"
        ]
    )

    # SLIDE 3: The Opportunity
    add_content_slide(
        prs,
        "The Opportunity",
        [
            "When volatility between related stocks diverges, it tends to snap back",
            "Semiconductor equipment makers vs chip designers",
            "Integrated oil majors vs refiners",
            "Temporary dislocations create trading opportunities"
        ],
        "Mean-reversion happens 5-7% of the time — we only trade when it's statistically extreme.",
        notes=[
            "Think of ASML and NVIDIA — both are in semiconductors, but ASML makes the machines that make the chips",
            "When fear hits the sector, they often move together, but sometimes one overreacts relative to the other",
            "The key insight is that we're not betting on direction — we're betting on the relationship normalizing"
        ]
    )

    # SLIDE 4: Basket Construction
    add_table_slide(
        prs,
        "Basket Construction",
        ["Pair", "Long Basket", "Short Basket"],
        [
            ["Semiconductors", "ASML, TSM, KLAC", "AMD, NVDA, AVGO"],
            ["Energy", "XOM, CVX, COP", "VLO, MPC, PSX"],
            ["Tech Broad vs Mega", "RSPT, SOXX", "QQQ, AAPL, META"],
            ["Staples vs Discretionary", "XLP", "XLY"]
        ],
        "Pairs selected for economic linkage — same sector, different volatility profiles.",
        notes=[
            "We chose these pairs because they have fundamental economic relationships, not just statistical correlation",
            "The long basket typically has lower volatility — equipment makers, integrated majors, broader ETFs",
            "Using baskets instead of single stocks reduces idiosyncratic risk from earnings surprises or company-specific news"
        ]
    )

    # SLIDE 5: Signal Generation
    add_two_column_slide(
        prs,
        "Signal Generation",
        "ENTRY RULES",
        [
            "Z-score > +2.0 → Short spread",
            "Z-score < -2.0 → Long spread",
            "VIX > 30 → Stay flat"
        ],
        "EXIT RULES",
        [
            "|Z-score| < 0.5 → Close position",
            "|Z-score| > 3.5 → Stop out",
            "Loss > 7% → Stop out"
        ],
        "We enter on extremes (2σ) and exit when spreads normalize or risk limits hit.",
        notes=[
            "The 2-sigma threshold is a balance — tighter means more trades with lower conviction, wider means fewer opportunities",
            "VIX filter is crucial — we learned that trading during crisis periods (like March 2020) destroys returns",
            "The 7% stop-loss was added after our backtest audit revealed we were letting losers run too long"
        ]
    )

    # SLIDE 6: ML Enhancement
    add_content_slide(
        prs,
        "Machine Learning Filter",
        [
            "Random Forest classifier filters raw Z-score signals",
            "8 features: z-score, momentum, vol_ratio, VIX, correlation",
            "Walk-forward validation with quarterly retraining",
            "30-day embargo between train and test periods"
        ],
        "ML acts as a quality filter — only take signals the model thinks will work.",
        notes=[
            "We tried 40+ features initially but found that simpler models with 8 core features performed better out-of-sample",
            "The 30-day embargo prevents data leakage — autocorrelation in returns can inflate backtest performance",
            "Random Forest was chosen for interpretability — we can see which features matter, unlike neural networks"
        ]
    )

    # SLIDE 7: Feature Importance
    add_chart_image_slide(
        prs,
        "What Drives the Model?",
        "feature_importance.png",
        "Z-score and momentum are the strongest predictors — simple features beat complex ones.",
        notes=[
            "Notice that the z-score itself is most important — the ML is essentially learning when z-score signals are reliable",
            "Momentum matters because mean-reversion strategies fail when there's a true regime shift happening",
            "VIX level helps the model avoid trading during extreme market stress when correlations break down"
        ]
    )

    # SLIDE 8: Backtest Setup
    add_content_slide(
        prs,
        "Backtest Methodology",
        [
            "Period: 2015-01-01 to 2024-12-31 (10 years)",
            "Position lagged 1 day (no look-ahead bias)",
            "Transaction costs: 5 bps per side (10 bps round-trip)",
            "Walk-forward quarterly retraining with 30-day embargo"
        ],
        "Backtest is clean — no look-ahead bias, realistic transaction costs.",
        notes=[
            "We actually found and fixed look-ahead bias during our audit — initial results looked much better before we fixed it",
            "10 bps round-trip is conservative for liquid stocks but accounts for market impact on larger positions",
            "Walk-forward testing is harder to implement but critical — static train/test splits overfit to specific market regimes"
        ]
    )

    # SLIDE 9: Cumulative Returns (ML vs Baseline vs SPY)
    add_chart_image_slide(
        prs,
        "Cumulative Returns: ML vs Baseline vs SPY",
        "cumulative_returns.png",
        "Three-way comparison shows ML improvement over baseline, both compared to SPY benchmark.",
        notes=[
            "The gray line is our baseline strategy, red is ML-enhanced — you can see ML reduces some drawdowns",
            "Semiconductors and Tech pairs show the most promise; Energy and Staples struggle to generate positive returns",
            "Remember, the goal isn't to beat SPY — it's to provide uncorrelated returns for portfolio diversification"
        ]
    )

    # SLIDE 10: Drawdown
    add_chart_image_slide(
        prs,
        "Drawdown Profile",
        "drawdown_chart.png",
        "Significant drawdowns require strong conviction and proper position sizing.",
        notes=[
            "50% drawdowns are painful — this is why we recommend only 5-10% portfolio allocation to this strategy",
            "Notice how drawdowns cluster around market stress periods like 2020 — mean-reversion fails when correlations spike",
            "The ML filter helps reduce some drawdowns but doesn't eliminate them — risk management is still essential"
        ]
    )

    # SLIDE 11: Performance Table
    add_table_slide(
        prs,
        "Performance: ML vs Baseline",
        ["Pair", "Baseline Sharpe", "ML Sharpe", "Baseline MaxDD", "ML MaxDD"],
        [
            ["Semiconductors", "-0.19", "+0.09", "-50%", "-30%"],
            ["Energy", "+0.01", "+0.08", "-50%", "-50%"],
            ["Tech vs Mega", "+0.18", "+0.09", "-26%", "-24%"],
            ["Staples vs Discr.", "-0.34", "-0.34", "-43%", "-35%"]
        ],
        "ML reduces drawdowns but doesn't consistently improve Sharpe — filtering helps risk, not return.",
        notes=[
            "Sharpe ratios near zero tell us this isn't a standalone alpha strategy — but that's okay for a hedge",
            "The ML improved Semiconductors significantly — from negative to positive Sharpe with 20% less drawdown",
            "Staples vs Discretionary is our worst performer — the relationship may be too noisy for mean-reversion"
        ]
    )

    # SLIDE 12: Rolling Sharpe
    add_chart_image_slide(
        prs,
        "Rolling Sharpe Ratio (252-Day)",
        "rolling_sharpe.png",
        "Sharpe fluctuates significantly — strategy has periods of strength and weakness.",
        notes=[
            "The wide swings between positive and negative Sharpe show this strategy requires patience and conviction",
            "You can see periods where the strategy works well (2018, 2022) and periods where it struggles (2020-2021)",
            "This volatility in performance is why we don't recommend this as a primary strategy — it's supplemental"
        ]
    )

    # SLIDE 13: Alpha Analysis
    add_table_slide(
        prs,
        "Alpha & Beta vs SPY",
        ["Pair", "Alpha", "Beta", "Corr w/ SPY"],
        [
            ["Semiconductors", "+1.2%", "-0.02", "-0.08"],
            ["Energy", "-2.1%", "+0.01", "+0.03"],
            ["Tech vs Mega", "+0.5%", "-0.03", "-0.12"],
            ["Staples vs Discr.", "-3.5%", "-0.01", "-0.05"],
            ["PORTFOLIO", "-0.89%", "-0.02", "-0.07"]
        ],
        "Near-zero beta confirms market neutrality — strategy returns are uncorrelated with SPY.",
        notes=[
            "Beta near zero across all pairs confirms the strategy is truly market-neutral — it won't move with SPY",
            "Negative correlation is actually desirable — it means this strategy can hedge equity exposure in a portfolio",
            "The -0.07 portfolio correlation means adding this to a 60/40 portfolio would reduce overall volatility"
        ]
    )

    # SLIDE 14: Monthly Heatmap
    add_chart_image_slide(
        prs,
        "Monthly Returns Heatmap",
        "monthly_heatmap.png",
        "No clear seasonality — returns are spread across different periods.",
        notes=[
            "We looked for monthly patterns like 'sell in May' but found no consistent seasonality in our returns",
            "The lack of seasonality is actually good — it means returns aren't driven by calendar effects that could disappear",
            "You can see the red cluster in early 2020 — that's the COVID crash where mean-reversion completely failed"
        ]
    )

    # SLIDE 15: Yearly vs SPY
    add_table_slide(
        prs,
        "Yearly Performance vs SPY",
        ["Year", "Strategy", "SPY", "Outperform?"],
        [
            ["2018", "+13%", "-4%", "Yes"],
            ["2022", "+16%", "-18%", "Yes"],
            ["2020", "-40%", "+18%", "No"],
            ["2021", "-25%", "+29%", "No"]
        ],
        "Strategy outperforms in down markets — potential use as a tail-risk hedge.",
        notes=[
            "This is the key insight: we outperform when SPY is negative (2018, 2022) and underperform in bull markets",
            "2020 was our worst year because COVID caused correlations to spike to 1 — everything fell together",
            "If you already own SPY, adding this strategy provides insurance in down years at the cost of bull market returns"
        ]
    )

    # SLIDE 16: Correlation Matrix
    add_chart_image_slide(
        prs,
        "Strategy Diversification",
        "correlation_matrix.png",
        "Low correlation between pairs (avg 0.03) means combining them reduces portfolio risk.",
        notes=[
            "Each pair trades independently — semiconductors don't predict energy, tech doesn't predict staples",
            "This low cross-correlation is powerful: combining all 4 pairs reduces volatility without reducing expected return",
            "In portfolio theory terms, we get diversification benefit from running multiple uncorrelated sub-strategies"
        ]
    )

    # SLIDE 17: Risk Summary
    add_table_slide(
        prs,
        "Risk Profile",
        ["Metric", "Value", "Interpretation"],
        [
            ["Max Drawdown", "-26% to -50%", "Significant capital risk"],
            ["VaR (95%)", "-1.0% to -1.7%", "Daily loss expectation"],
            ["CVaR (95%)", "-1.6% to -2.9%", "Tail risk worst days"],
            ["Win Rate", "48-52%", "Below coin flip"]
        ],
        "This is a low win-rate, high-variance strategy — position sizing is critical.",
        notes=[
            "50% max drawdown means you need strong conviction — most investors would abandon the strategy mid-drawdown",
            "VaR and CVaR tell us on a bad day (5% worst days), we lose 1-3% which is manageable with proper sizing",
            "Win rate near 50% means profits come from winners being bigger than losers, not from winning more often"
        ]
    )

    # SLIDE 18: Lessons Learned
    add_two_column_slide(
        prs,
        "Fixes & Lessons Learned",
        "WHAT WE FIXED",
        [
            "✓ Removed look-ahead bias",
            "✓ Added 30-day train/test embargo",
            "✓ Corrected transaction costs",
            "✓ Implemented 3 stop-loss rules"
        ],
        "WHAT WE LEARNED",
        [
            "• 8 features beat 40 features",
            "• Spread mean-reverts only 5-7%",
            "• VIX regime matters significantly",
            "• Alpha is hard to find"
        ],
        "Proper backtesting revealed our initial results were inflated — honesty improved the strategy.",
        notes=[
            "Our initial backtest showed 2+ Sharpe ratio — after fixing biases, we got near-zero. That's a humbling lesson.",
            "Look-ahead bias is the #1 killer of academic trading strategies — we used next-day signals accidentally",
            "The fact that we found and fixed these issues is actually the most valuable part of this project"
        ]
    )

    # SLIDE 19: Conclusion
    add_two_column_slide(
        prs,
        "Investment Thesis",
        "WHAT WORKED",
        [
            "✓ Market-neutral (β ≈ 0)",
            "✓ Low SPY correlation",
            "✓ Stop-losses limit tail risk",
            "✓ Outperforms in down markets"
        ],
        "WHAT DIDN'T",
        [
            "✗ No consistent alpha",
            "✗ Large drawdowns (26-50%)",
            "✗ Underperforms in bull markets",
            "✗ Low win rate (~50%)"
        ],
        "Best suited as a PORTFOLIO HEDGE — allocate 5-10% for downside protection.",
        notes=[
            "We're not claiming this is an alpha machine — we're being honest about what it is: a diversification tool",
            "The 5-10% allocation recommendation comes from balancing hedge benefit against opportunity cost in bull markets",
            "For a family office or pension fund, this type of uncorrelated strategy has real value even at low Sharpe"
        ]
    )

    # SLIDE 20: Questions
    add_questions_slide(
        prs,
        notes=[
            "Thank you for your attention — all code is available on GitHub for you to reproduce our results",
            "We welcome questions about methodology, the ML approach, or how this could fit in a real portfolio",
            "This project taught us more about proper backtesting than any textbook — happy to discuss the technical details"
        ]
    )

    # Save
    output_path = "/Users/akbarpathan/Desktop/Dev/QuantHFStrat/FE571_Final_Presentation_v2.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
