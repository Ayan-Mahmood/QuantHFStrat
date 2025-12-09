#!/usr/bin/env python3
"""
Create professional FinTech-style pitch deck for FE571 Final Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Stevens Institute Colors
STEVENS_RED = RGBColor(0x9D, 0x15, 0x35)
STEVENS_GRAY = RGBColor(0x94, 0x95, 0x94)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

def set_slide_background(slide, color=WHITE):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle, authors):
    """Add title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = STEVENS_RED
    p.alignment = PP_ALIGN.CENTER

    # Authors
    auth_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    tf = auth_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FE571 | Professor Anshul Sharma | Group 7"
    p.font.size = Pt(16)
    p.font.color.rgb = STEVENS_GRAY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = authors
    p.font.size = Pt(14)
    p.font.color.rgb = STEVENS_GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, one_liner=None):
    """Add a content slide with title, bullets, and optional one-liner"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Red accent line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STEVENS_RED
    line.line.fill.background()

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    # One-liner at bottom
    if one_liner:
        liner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
        tf = liner_box.text_frame
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = STEVENS_GRAY

    return slide

def add_table_slide(prs, title, headers, rows, one_liner=None):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Red accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STEVENS_RED
    line.line.fill.background()

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(9)
    table_height = Inches(0.4 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.4), table_width, table_height).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = STEVENS_RED
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = DARK_GRAY
            para.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # One-liner
    if one_liner:
        liner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
        tf = liner_box.text_frame
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = STEVENS_GRAY

    return slide

def add_chart_placeholder_slide(prs, title, chart_name, one_liner=None):
    """Add a slide with placeholder for chart"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Red accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STEVENS_RED
    line.line.fill.background()

    # Chart placeholder box
    chart_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.4), Inches(8.5), Inches(3.5))
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    chart_box.line.color.rgb = STEVENS_GRAY

    # Placeholder text
    text_box = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(6), Inches(1))
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"[INSERT: {chart_name}]"
    p.font.size = Pt(18)
    p.font.color.rgb = STEVENS_GRAY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Export from v5 notebook"
    p.font.size = Pt(12)
    p.font.color.rgb = STEVENS_GRAY
    p.alignment = PP_ALIGN.CENTER

    # One-liner
    if one_liner:
        liner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
        tf = liner_box.text_frame
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = STEVENS_GRAY

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, one_liner=None):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    # Red accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(1.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = STEVENS_RED
    line.line.fill.background()

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = STEVENS_RED

    # Left column items
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(3))
    tf = left_box.text_frame
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = STEVENS_RED

    # Right column items
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4.2), Inches(3))
    tf = right_box.text_frame
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # One-liner
    if one_liner:
        liner_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
        tf = liner_box.text_frame
        p = tf.paragraphs[0]
        p.text = one_liner
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = STEVENS_GRAY

    return slide

def add_questions_slide(prs):
    """Add questions slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    # Questions text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    # GitHub link
    link_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
    tf = link_box.text_frame
    p = tf.paragraphs[0]
    p.text = "github.com/Ayan-Mahmood/QuantHFStrat"
    p.font.size = Pt(16)
    p.font.color.rgb = STEVENS_RED
    p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

    # SLIDE 1: Title
    add_title_slide(
        prs,
        "Statistical Basket Pairs Trading Strategy",
        "Volatility Dispersion Mean-Reversion",
        "Scott Henriquez, Nakul Jadeja, Ayan Mahmood, Akbar Pathan"
    )

    # SLIDE 2: Executive Summary
    add_content_slide(
        prs,
        "Executive Summary",
        [
            "Long/short volatility spread strategy across 4 sector pairs",
            "Market-neutral design with low SPY correlation (-0.07)",
            "ML filter tested to improve signal quality",
            "Result: Solid framework, inconsistent alpha — best as portfolio hedge"
        ]
    )

    # SLIDE 3: The Opportunity
    add_content_slide(
        prs,
        "The Opportunity",
        [
            "When volatility between related stocks diverges, it tends to snap back",
            "Semiconductor equipment makers vs chip designers",
            "Integrated oil majors vs refiners",
            "Temporary dislocations create trading opportunities"
        ],
        "Mean-reversion happens 5-7% of the time — we only trade when it's statistically extreme."
    )

    # SLIDE 4: Basket Construction
    add_table_slide(
        prs,
        "Basket Construction",
        ["Pair", "Long Basket", "Short Basket"],
        [
            ["Semiconductors", "ASML, TSM, KLAC", "AMD, NVDA, AVGO"],
            ["Energy", "XOM, CVX, COP", "VLO, MPC, PSX"],
            ["Tech Broad vs Mega", "RSPT, SOXX", "QQQ, AAPL, META"],
            ["Staples vs Discretionary", "XLP", "XLY"]
        ],
        "Pairs selected for economic linkage — same sector, different volatility profiles."
    )

    # SLIDE 5: Signal Generation
    add_two_column_slide(
        prs,
        "Signal Generation",
        "ENTRY RULES",
        [
            "Z-score > +2.0 → Short spread",
            "Z-score < -2.0 → Long spread",
            "VIX > 30 → Stay flat"
        ],
        "EXIT RULES",
        [
            "|Z-score| < 0.5 → Close position",
            "|Z-score| > 3.5 → Stop out",
            "Loss > 7% → Stop out"
        ],
        "We enter on extremes (2σ) and exit when spreads normalize or risk limits hit."
    )

    # SLIDE 6: ML Enhancement
    add_content_slide(
        prs,
        "Machine Learning Filter",
        [
            "Random Forest classifier filters raw Z-score signals",
            "8 features: z-score, momentum, vol_ratio, VIX, correlation",
            "Walk-forward validation with quarterly retraining",
            "30-day embargo between train and test periods"
        ],
        "ML acts as a quality filter — only take signals the model thinks will work."
    )

    # SLIDE 7: Feature Importance
    add_chart_placeholder_slide(
        prs,
        "What Drives the Model?",
        "feature_importance.png",
        "Z-score and momentum are the strongest predictors — simple features beat complex ones."
    )

    # SLIDE 8: Backtest Setup
    add_content_slide(
        prs,
        "Backtest Methodology",
        [
            "Period: 2015-01-01 to 2024-12-31 (10 years)",
            "Position lagged 1 day (no look-ahead bias)",
            "Transaction costs: 5 bps per side (10 bps round-trip)",
            "Walk-forward quarterly retraining with 30-day embargo"
        ],
        "Backtest is clean — no look-ahead bias, realistic transaction costs."
    )

    # SLIDE 9: Cumulative Returns
    add_chart_placeholder_slide(
        prs,
        "Cumulative Returns",
        "cumulative_returns.png",
        "Mixed results across pairs — Semiconductors and Tech show promise, Energy and Staples struggle."
    )

    # SLIDE 10: Drawdown
    add_chart_placeholder_slide(
        prs,
        "Drawdown Profile",
        "drawdown_chart.png",
        "Significant drawdowns require strong conviction and proper position sizing."
    )

    # SLIDE 11: Performance Table
    add_table_slide(
        prs,
        "Performance: ML vs Baseline",
        ["Pair", "Baseline Sharpe", "ML Sharpe", "Baseline MaxDD", "ML MaxDD"],
        [
            ["Semiconductors", "-0.19", "+0.09", "-50%", "-30%"],
            ["Energy", "+0.01", "+0.08", "-50%", "-50%"],
            ["Tech vs Mega", "+0.18", "+0.09", "-26%", "-24%"],
            ["Staples vs Discr.", "-0.34", "-0.34", "-43%", "-35%"]
        ],
        "ML reduces drawdowns but doesn't consistently improve Sharpe — filtering helps risk, not return."
    )

    # SLIDE 12: Rolling Sharpe
    add_chart_placeholder_slide(
        prs,
        "Rolling Sharpe Ratio (252-Day)",
        "rolling_sharpe.png",
        "Sharpe fluctuates significantly — strategy has periods of strength and weakness."
    )

    # SLIDE 13: Alpha Analysis
    add_table_slide(
        prs,
        "Alpha & Beta vs SPY",
        ["Pair", "Alpha", "Beta", "Corr w/ SPY"],
        [
            ["Semiconductors", "+1.2%", "-0.02", "-0.08"],
            ["Energy", "-2.1%", "+0.01", "+0.03"],
            ["Tech vs Mega", "+0.5%", "-0.03", "-0.12"],
            ["Staples vs Discr.", "-3.5%", "-0.01", "-0.05"],
            ["PORTFOLIO", "-0.89%", "-0.02", "-0.07"]
        ],
        "Near-zero beta confirms market neutrality — strategy returns are uncorrelated with SPY."
    )

    # SLIDE 14: Monthly Heatmap
    add_chart_placeholder_slide(
        prs,
        "Monthly Returns Heatmap",
        "monthly_heatmap.png",
        "No clear seasonality — returns are spread across different periods."
    )

    # SLIDE 15: Yearly vs SPY
    add_table_slide(
        prs,
        "Yearly Performance vs SPY",
        ["Year", "Strategy", "SPY", "Outperform?"],
        [
            ["2018", "+13%", "-4%", "Yes"],
            ["2022", "+16%", "-18%", "Yes"],
            ["2020", "-40%", "+18%", "No"],
            ["2021", "-25%", "+29%", "No"]
        ],
        "Strategy outperforms in down markets — potential use as a tail-risk hedge."
    )

    # SLIDE 16: Correlation Matrix
    add_chart_placeholder_slide(
        prs,
        "Strategy Diversification",
        "correlation_matrix.png",
        "Low correlation between pairs (avg 0.03) means combining them reduces portfolio risk."
    )

    # SLIDE 17: Risk Summary
    add_table_slide(
        prs,
        "Risk Profile",
        ["Metric", "Value", "Interpretation"],
        [
            ["Max Drawdown", "-26% to -50%", "Significant capital risk"],
            ["VaR (95%)", "-1.0% to -1.7%", "Daily loss expectation"],
            ["CVaR (95%)", "-1.6% to -2.9%", "Tail risk worst days"],
            ["Win Rate", "48-52%", "Below coin flip"]
        ],
        "This is a low win-rate, high-variance strategy — position sizing is critical."
    )

    # SLIDE 18: Lessons Learned
    add_two_column_slide(
        prs,
        "Fixes & Lessons Learned",
        "WHAT WE FIXED",
        [
            "✓ Removed look-ahead bias",
            "✓ Added 30-day train/test embargo",
            "✓ Corrected transaction costs",
            "✓ Implemented 3 stop-loss rules"
        ],
        "WHAT WE LEARNED",
        [
            "• 8 features beat 40 features",
            "• Spread mean-reverts only 5-7%",
            "• VIX regime matters significantly",
            "• Alpha is hard to find"
        ],
        "Proper backtesting revealed our initial results were inflated — honesty improved the strategy."
    )

    # SLIDE 19: Conclusion
    add_two_column_slide(
        prs,
        "Investment Thesis",
        "WHAT WORKED",
        [
            "✓ Market-neutral (β ≈ 0)",
            "✓ Low SPY correlation",
            "✓ Stop-losses limit tail risk",
            "✓ Outperforms in down markets"
        ],
        "WHAT DIDN'T",
        [
            "✗ No consistent alpha",
            "✗ Large drawdowns (26-50%)",
            "✗ Underperforms in bull markets",
            "✗ Low win rate (~50%)"
        ],
        "Best suited as a PORTFOLIO HEDGE — allocate 5-10% for downside protection."
    )

    # SLIDE 20: Questions
    add_questions_slide(prs)

    # Save
    output_path = "/Users/akbarpathan/Desktop/Dev/QuantHFStrat/FE571_Final_Presentation_v2.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
